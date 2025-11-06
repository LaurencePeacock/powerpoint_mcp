from mcp.server import FastMCP
import logging
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from mcp.server.fastmcp import Context

from SessionManager import SessionManager
from utils.clean_slide_name import clean_slide_name
from utils.validate_table_data import validate_table_data
from errors.TableDataValidationException import TableDataValidationException

logger = logging.getLogger(__name__)


def add_data_to_table(
        table,
        table_data: dict,
        num_cols: int,
        num_rows: int,
) -> None:
    logger.info(f'---- Attempting to add data to table.')
    # get columns
    for i in range(num_cols):
        cell_x = 0
        cell_y = i
        cell = table.cell(cell_x, cell_y)
        cell.text = str(table_data['columns'][i])
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # for each column, access the datasets
    for j in range(0, num_rows):
        # for each dataset, access each value
        for k in range(0, num_cols):
            cell_x = j + 1
            cell_y = k
            cell = table.cell(cell_x, cell_y)
            cell.text = str(table_data['values'][j][k])
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def register_table_tools(
        pp_app: FastMCP,
        session_manager: SessionManager
):
    @pp_app.tool()
    def table_handler(table_data: dict):
        """
        Receives a dictionary of data from the root agent and checks if it meets the data structure requirements
        for table data.
        :param table_data: a dictionary of data in json format
        :return: a dictionary with a success or failure status. If successful, the validated data is also included.
        If unsuccessful, an error message is included.

        Example of Successful Return dictionary:
        {
            "status": "success",
            "table":{
                "columns": ["colA", "colB", "colC"],
                "values": [
                    ['a', 'b', 'c'],
                    ['d', 'e', 'f'],
                    ['g', 'h', 'i'],
                    ['j', 'k', 'l']
                ]
            }
        }
        Example of Failure return dictionary:
        {
            "status": "failure",
            "error_message": "Unable to convert provided data into table_data dictionary."
        }
        """
        logger.info(f"---- Table Handler called with data: {table_data}")
        try:
            validated, message = validate_table_data(table_data)
            if validated:
                return {
                    "status": "success",
                    "table_data": table_data
                }
            else:
                return {
                    "status": "failure",
                    "error_message": message
                }
        except TableDataValidationException as e:
            logger.error(f"---- Unable to convert json into table data format: {e}")
            return {
                "status": "failure",
                "error_message": f"Error Validating Table Data. Unable to validate table data: {e}"
            }



    @pp_app.tool()
    def add_table_to_slide(
            presentation_filename: str,
            table_data: dict,
            slide_name: str = None,
            slide_index: int = None
    )-> dict:
        """
        IMPORTANT - only use the slide_name or slide_index provided to find the slide. DO NOT USE the get_presentation_summary tool
        IMPORTANT - follow the tool code precisely. Do not attempt to work out your own plan.

        This tool finds a slide by its index or name and adds a table to the first available content placeholder.
        This only modifies the slide in memory. Call 'save_presentation' to persist changes.

        :param presentation_filename: the filename of the presentation file
        :param table_data: The json data which will be used for the table
        :param slide_name: Optional. The name of the slide to add text to, E.G 'title' or 'main'
        :param slide_index: Optional. The index of the slide to add text to (0-based).

        CRITICAL INSTRUCTION: The table_data you receive will already have been validated by the table_handler method.
        table_data MUST BE json structured with keys for 'columns' and 'values'.
        The "columns" dictionary MUST contain a list of strings or numbers
        The "values" dictionary MUST contain a multidimensional list of lists containing strings or numbers
        table_data Example:
             {
                "columns": ["colA", "colB", "colC"],
                "values": [
                    ['a', 'b', 'c'],
                    ['d', 'e', 'f'],
                    ['g', 'h', 'i'],
                    ['j', 'k', 'l']
                ]
            }
        """
        # add a table to the slide at index 1. Here is the table data: "{"columns": ["colA", "colB", "colC"],"values": [['a', 'b', 'c'], ['d', 'e', 'f']]}"
        logger.info(f"---- Attempting to add table to slide")
        presentation = session_manager.get_presentation(presentation_filename)
        if not presentation:
            return {
                "status": "failure",
                "message": "Presentation not found for the current session in add_new_slide_tool."
            }
        logger.info(
            f"---- Tool add_table_to_slide called with arguments: chart_data: {table_data}, slide_name: {slide_name}, slide_index: {slide_index}")
        try:
            if slide_index and slide_index >= len(presentation.slides):
                return f"Error: Invalid slide index. The presentation only has {len(presentation.slides)} slides."
        except Exception as e:
            return {
                "status": "failure",
                "error": f"{e}"
            }

        logger.info(f"---- Provided with slide name: {slide_name} and / or index {slide_index}")
        if slide_name:
            logger.info(f"---- Using the slide name to find the slide in the presentation")
            cleaned_name = clean_slide_name(slide_name)
            for named_slide in presentation.slides:
                if named_slide.name == slide_name:
                    slide_to_edit = named_slide
        else:
            logger.info(f"---- slide name was not provided. Using Index to find the slide in the presentation")
            slide_to_edit = presentation.slides[slide_index]

        if not slide_to_edit:
            return f"Error: Unable to find slide to add table to."

        # Find the first empty placeholder suitable for a table
        logger.info(f"---- finding placeholders on slide: {slide_to_edit.name}")
        logger.info(f"---- There are {len(slide_to_edit.placeholders)} placeholders on slide")
        table_already_on_slide = False
        for shape in slide_to_edit.placeholders:
            logger.info(f"---- Placeholder Type: {shape.placeholder_format.type}")
            if shape.placeholder_format.type.name in ('TABLE',):
                if shape.has_table:
                    logger.info(f"---- Table placeholder found on {slide_name} but it already contains table data")
                    table_already_on_slide = True
                    continue
                logger.info(f"---- found table placeholders on slide: {slide_to_edit.name}")
                table_placeholder = shape
                logger.info(f"---- constructing table object")

                num_rows = len(table_data["values"])
                num_cols = len(table_data["columns"])
                # rows count includes column headers so add one to length
                table_shape = table_placeholder.insert_table(rows=num_rows + 1, cols=num_cols)
                # insert_table returns a placeholder with the table stored in the table property. Hence to edit the actual table:
                table_to_edit = table_shape.table

                try:
                    add_data_to_table(
                        table_to_edit,
                        table_data,
                        num_cols,
                        num_rows
                    )
                except Exception as e:
                    logger.error(f'Error adding table data to table object on slide: {slide_to_edit.name}. Error: {e}')
                    return {
                        "status": "failure",
                        "message": f"add_data_to_table raised an error adding table data to slide: {slide_to_edit.name}"
                    }

                return {
                    "status": "success",
                    "message": f"Successfully added table to slide {slide_to_edit.name}, index: {slide_index}. Remember to save."
                }

        if table_already_on_slide:
            logger.warning(
                f"---- No empty table placeholder was found on slide {slide_to_edit.name}, index: {slide_index}. Please add another slide")
            return {
                "status": "failure",
                "message": f"Adding table to slide slide {slide_to_edit.name}, index: {slide_index}. because a populated table already exists on the slide.",
                "user_suggestion": "Suggest user adds a new slide with a Table Layout."
            }

        return {
            "status": "failure",
            "message": f"Adding table to slide slide {slide_to_edit.name}, index: {slide_index}. because no table placeholder was found on the slide. Please add another slide with a Table Layout.",
            "user_suggestion": "Suggest user adds a new slide with a Table Layout."
        }