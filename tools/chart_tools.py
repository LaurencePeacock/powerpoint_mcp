from mcp.server import FastMCP
import logging

from mcp.server.fastmcp import Context
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Pt
from SessionManager import SessionManager
from utils.clean_slide_name import clean_slide_name
from utils.validate_chart_data import validate_chart_data
from errors.ChartDataConverterException import ChartDataConverterException

logger = logging.getLogger(__name__)

def register_chart_tools(
        pp_app: FastMCP,
        session_manager: SessionManager
):
    @pp_app.tool()
    def chart_handler(chart_data: dict) -> dict:
        """
        Receives a dictionary of data from the root agent and attempts to convert it into a chart data format with categories and series
        :param chart_data: a dictionary of data in json format
        :return: a dictionary with a success or failure status. If successful, the converted data is also included.
        If unsuccessful, an error message is included.

        Example of Successful Return dictionary:
        {
            "status": "success",
            "chart_data":{
                "categories": ["2025-01-01", "2025-02-01", "2025-03-01", "2025-04-01"],
                "series" : [
                    {
                        "name": "total_sales",
                        "values": [6, 10, 4, 8]
                    }
                ]
            }
        }
        Example of Failure return dictionary:
        {
            "status": "failure",
            "error_message": "Unable to convert provided data into chart_data dictionary."
        }
        """

        logger.info(f"---- Chart Handler called with data: {chart_data}")
        logger.info(f"---- Attempting to turn json data into chart data ")
        try:
            validated, message = validate_chart_data(chart_data)
            if validated:
                return {
                    "status": "success",
                    "chart_data": chart_data
                }
            else:
                return {
                    "status": "failure",
                    "error_message": message
                }
        except ChartDataConverterException as e:
            logger.error(f"---- Unable to convert json into chart_data format: {e}")
            return {
                "status": "failure",
                "error_message": f"Error Validating Chart Data. Unable to validate chart_data: {e}"
            }

    @pp_app.tool()
    def add_chart_to_slide(
            presentation_filename: str,
            chart_data: dict,
            chart_title: str = None,
            chart_has_legend: bool = False,
            category_axis_title: str = None,
            value_axis_title: str = None,
            slide_name: str = None,
            slide_index: int = None
    ) -> dict:
        """
        IMPORTANT - only use the slide_name or slide_index provided to find the slide. DO NOT USE the get_presentation_summary tool
        IMPORTANT - follow the tool code precisely. Do not attempt to work out your own plan.

        This tool finds a slide by its index or name and adds a chart to the first available content placeholder.

        This only modifies the slide in memory. Call 'save_presentation' to persist changes.

        :param presentation_filename: the filename of the presentation file
        :param chart_data: The json data which will be used for the chart
        :param chart_title: the title of the chart. If this is not provided by the user, you should generate a short title based on the chart_data
        :param chart_has_legend: boolean indicating if this chart has a legend or not. By default this is set to True. If the user specifies that the chart
        should not have a legend, this must be set to False.
        :param category_axis_title: string. This must be provided by the Agent based on the data included in  chart_data["categories"]
        :param value_axis_title: string. This must be provided by the Agent based on the data included in chart_data["series"]. Include an appropriate unit of measurement.
        For example, if the name of the series is 'total spend', do 'total spend (£)'. If the name of the series concerns a percentage, do 'example tile (%)'
        :param slide_name: Optional. The name of the slide to add text to, E.G 'title' or 'main'
        :param slide_index: Optional. The index of the slide to add text to (0-based).

        :return: a dictionary indicating the success or failure of the tool

        CRITICAL INSTRUCTION: chart_data MUST BE json structured with keys for 'categories' and 'series' data.
        The "series" dictionary MUST HAVE "name" and "values" keys.
        chart_data Example:
             {
                "categories": ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
                "series": [
                        {"name": "Total Sales", "values": [6, 14, 16, 23, 28, 8]}
                    ]
             }
        """
        logger.info(f"---- Attempting to add chart to slide")
        presentation = session_manager.get_presentation(presentation_filename)
        if not presentation:
            return {
                "status": "failure",
                "message": "Presentation not found for the current session in add_new_slide_tool."
            }
        logger.info(
            f"---- Tool called with arguments: chart categories: {chart_data}, slide_name: {slide_name}, slide_index: {slide_index}")
        if slide_index and slide_index >= len(presentation.slides):
            return f"Error: Invalid slide index. The presentation only has {len(presentation.slides)} slides."
        logger.info(f"---- Provided with slide name: {slide_name} and / or index {slide_index}")
        if slide_name:
            logger.info(f"---- Using the slide name to find the slide in the presentation")
            cleaned_name = clean_slide_name(slide_name)
            for named_slide in presentation.slides:
                if named_slide.name == cleaned_name:
                    slide_to_edit = named_slide
        else:
            logger.info(f"---- slide name was not provided. Using Index to find the slide in the presentation")
            slide_to_edit = presentation.slides[slide_index]

        if not slide_to_edit:
            return f"Error: Unable to find slide to add chart to."

        # Find the first empty placeholder suitable for a chart
        chart_placeholder = None
        logger.info(f"---- finding placeholders on slide: {slide_to_edit.name}")
        logger.info(f"---- There are {len(slide_to_edit.placeholders)} placeholders on slide")
        placeholders = slide_to_edit.placeholders
        chart_already_on_slide = False
        for shape in slide_to_edit.placeholders:
            logger.info(f"---- Placeholder Type: {shape.placeholder_format.type}")
            if shape.placeholder_format.type.name in ('CHART',):
                if shape.has_chart:
                    logger.info(f"---- Chart placeholder found on {slide_name} but it already contains a chart")
                    chart_already_on_slide = True
                    continue
                logger.info(f"---- found chart placeholders on slide: {slide_to_edit.name}")
                chart_placeholder = shape
                logger.info(f"---- constructing chart object")

                data = CategoryChartData()
                data.categories = chart_data["categories"]
                for series in chart_data["series"]:
                    data.add_series(series.get('name', 'Unnamed Series'), series.get('values', []))

                chart_graphic_frame = chart_placeholder.insert_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    data
                )

                chart = chart_graphic_frame.chart
                if chart_title:
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_title
                    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

                # set category axis titles and font size
                category_axis = chart.category_axis
                category_axis.axis_title.text_frame.text = category_axis_title
                category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(13)

                # set font size of axis data labels
                category_axis.tick_labels.font.size = Pt(11)
                category_axis.visible = True

                # set value axis titles and font size
                value_axis = chart.value_axis
                value_axis.axis_title.text_frame.text = value_axis_title
                value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(13)

                # set font size of axis data labels
                value_axis.tick_labels.font.size = Pt(11)
                value_axis.visible = True


                if chart_has_legend:
                    chart.has_legend = True
                    # slide.name uses user_friendly_name property from slide_layouts_full.yaml
                    if slide_to_edit.name in (
                            'Chart and Text - grey background',
                            'Chart Right Hand with Text on Left',
                            'Text Content with Left Hand Graph',
                    ):
                        legend_position = XL_LEGEND_POSITION.BOTTOM
                    else:
                        legend_position = XL_LEGEND_POSITION.RIGHT
                    chart.legend.position = legend_position
                    chart.legend.font.size = Pt(11)
                    chart.legend.include_in_layout = False


                chart_already_on_slide = False
                return {
                    "status": "success",
                    "message": f"Successfully added chart to slide {slide_to_edit.name}, index: {slide_index}. Remember to save."
                }
        if chart_already_on_slide:
            return {
                "status": "failure",
                "message": 'chart already on slide error message'
            }
        logger.warning(
            f"---- No empty chart placeholder was found on slide {slide_to_edit.name}, index: {slide_index}. Please add another slide")
        return {
            "status": "failure",
            "message": f"Error: This is because no chart placeholder was found on the slide. Please add another slide with a Chart Layout."
        }